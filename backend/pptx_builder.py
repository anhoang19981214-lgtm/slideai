import io
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

THEMES = {
    "purple": {"bg": (10, 6, 24),   "bg2": (21, 13, 53),  "accent": (167,139,250), "accent2": (192,132,252), "text": (237,233,254), "muted": (196,181,253)},
    "pink":   {"bg": (19, 0,  8),   "bg2": (37,  0, 21),  "accent": (244,114,182), "accent2": (249,168,212), "text": (252,231,243), "muted": (249,168,212)},
    "blue":   {"bg": (0,  8,  15),  "bg2": (0,  20, 40),  "accent": (96, 165,250), "accent2": (147,197,253), "text": (219,234,254), "muted": (147,197,253)},
    "green":  {"bg": (0,  16, 10),  "bg2": (0,  32, 24),  "accent": (52, 211,153), "accent2": (110,231,183), "text": (209,250,229), "muted": (110,231,183)},
    "sunset": {"bg": (15, 5,  0),   "bg2": (30, 8,  0),   "accent": (251,146, 60), "accent2": (251,191, 36), "text": (255,247,237), "muted": (252,211, 77)},
}

W = Inches(13.33)
H = Inches(7.5)
MARGIN = Inches(0.6)


def _rgb(t): return RGBColor(*t)


def _set_bg(slide, color):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = _rgb(color)


def _rect(slide, left, top, width, height, color, alpha=None):
    shape = slide.shapes.add_shape(1, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = _rgb(color)
    shape.line.fill.background()
    return shape


def _textbox(slide, text, left, top, width, height, size, bold=False, color=None, align=PP_ALIGN.LEFT, italic=False, wrap=True):
    if not text:
        return
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = wrap
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.italic = italic
    if color:
        run.font.color.rgb = _rgb(color)
    return txBox


def _cover_slide(slide, s, colors):
    _set_bg(slide, colors["bg"])

    # Top accent stripe
    _rect(slide, 0, 0, W, Inches(0.06), colors["accent"])

    # Decorative circle (right side, large)
    cx = Inches(9.5)
    size = Inches(5)
    sh = slide.shapes.add_shape(9, cx, Inches(1.2), size, size)  # oval
    mid = tuple(int(colors["bg"][i] * 0.7 + colors["accent"][i] * 0.3) for i in range(3))
    sh.fill.solid()
    sh.fill.fore_color.rgb = _rgb(mid)
    sh.line.fill.background()

    # Inner smaller circle
    sh2 = slide.shapes.add_shape(9, cx + Inches(1), Inches(2), Inches(3), Inches(3))
    mid2 = tuple(int(colors["bg"][i] * 0.5 + colors["accent"][i] * 0.5) for i in range(3))
    sh2.fill.solid()
    sh2.fill.fore_color.rgb = _rgb(mid2)
    sh2.line.fill.background()

    # Accent tag line
    _rect(slide, MARGIN, Inches(1.8), Inches(0.05), Inches(0.5), colors["accent"])

    # Title
    _textbox(slide, s.get("title", ""), MARGIN + Inches(0.2), Inches(1.7), Inches(8.5), Inches(1.8),
             44, bold=True, color=colors["text"], align=PP_ALIGN.LEFT)

    # Divider bar
    _rect(slide, MARGIN, Inches(3.65), Inches(1.2), Inches(0.05), colors["accent"])

    # Subtitle
    _textbox(slide, s.get("subtitle", ""), MARGIN, Inches(3.9), Inches(8.5), Inches(1.2),
             22, bold=False, color=colors["muted"], align=PP_ALIGN.LEFT)

    # Bottom bar
    _rect(slide, 0, H - Inches(0.06), W, Inches(0.06), colors["accent2"])

    # Slide number (bottom right)
    _textbox(slide, "01", Inches(11.5), Inches(6.8), Inches(1.5), Inches(0.5),
             14, bold=True, color=colors["muted"], align=PP_ALIGN.RIGHT)


def _content_slide(slide, s, colors, idx):
    _set_bg(slide, colors["bg"])

    # Left accent bar
    _rect(slide, 0, 0, Inches(0.08), H, colors["accent"])

    # Top stripe
    _rect(slide, 0, 0, W, Inches(0.04), colors["accent2"])

    # Slide number (watermark, large faint)
    num_str = str(idx).zfill(2)

    # Title area background
    _rect(slide, Inches(0.08), 0, W - Inches(0.08), Inches(1.5), colors["bg2"])

    # Slide label
    _textbox(slide, f"SLIDE {num_str}", Inches(0.3), Inches(0.15), Inches(3), Inches(0.4),
             9, bold=True, color=colors["accent"], align=PP_ALIGN.LEFT)

    # Title
    _textbox(slide, s.get("title", ""), Inches(0.3), Inches(0.45), Inches(12.5), Inches(0.95),
             28, bold=True, color=colors["text"], align=PP_ALIGN.LEFT)

    # Title underline
    _rect(slide, Inches(0.3), Inches(1.38), Inches(0.8), Inches(0.04), colors["accent"])

    # Bullets
    bullets = s.get("bullets", [])
    cols = 2 if len(bullets) >= 4 else 1
    col_w = (Inches(12.5) - Inches(0.2)) / cols
    row_h = min(Inches(1.05), (H - Inches(1.7)) / max(1, (len(bullets) + cols - 1) // cols))

    for i, bullet in enumerate(bullets):
        col = i % cols
        row = i // cols
        bx = Inches(0.3) + col * (col_w + Inches(0.2))
        by = Inches(1.65) + row * (row_h + Inches(0.1))

        # Card background
        _rect(slide, bx, by, col_w, row_h, colors["bg2"])
        # Left accent
        _rect(slide, bx, by, Inches(0.05), row_h, colors["accent"])
        # Number circle
        _textbox(slide, str(i + 1).zfill(2), bx + Inches(0.12), by + Inches(0.15), Inches(0.5), Inches(0.5),
                 10, bold=True, color=colors["accent"], align=PP_ALIGN.CENTER)
        # Bullet text
        _textbox(slide, bullet, bx + Inches(0.68), by + Inches(0.1), col_w - Inches(0.8), row_h - Inches(0.2),
                 14, bold=False, color=colors["text"], align=PP_ALIGN.LEFT)

    # Bottom bar
    _rect(slide, 0, H - Inches(0.04), W, Inches(0.04), colors["accent"])

    # Slide number bottom right
    _textbox(slide, num_str, Inches(12), Inches(6.9), Inches(1.2), Inches(0.5),
             13, bold=True, color=colors["muted"], align=PP_ALIGN.RIGHT)


def _conclusion_slide(slide, s, colors, idx):
    _set_bg(slide, colors["bg"])

    # Top accent stripe
    _rect(slide, 0, 0, W, Inches(0.06), colors["accent"])

    # Center glow circle (simulated with oval)
    for size_i, alpha_factor in [(5, 0.15), (3.5, 0.25), (2, 0.4)]:
        ox = (W - Inches(size_i)) / 2
        oy = (H - Inches(size_i)) / 2
        sh = slide.shapes.add_shape(9, ox, oy, Inches(size_i), Inches(size_i))
        blend = tuple(int(colors["bg"][j] * (1 - alpha_factor) + colors["accent"][j] * alpha_factor) for j in range(3))
        sh.fill.solid()
        sh.fill.fore_color.rgb = _rgb(blend)
        sh.line.fill.background()

    # Icon circle
    icon_size = Inches(0.7)
    ix = (W - icon_size) / 2
    sh_i = slide.shapes.add_shape(9, ix, Inches(2), icon_size, icon_size)
    sh_i.fill.solid()
    sh_i.fill.fore_color.rgb = _rgb(colors["accent"])
    sh_i.line.fill.background()

    # Title
    _textbox(slide, s.get("title", ""), MARGIN, Inches(2.9), W - MARGIN * 2, Inches(1.3),
             36, bold=True, color=colors["text"], align=PP_ALIGN.CENTER)

    # Divider
    div_w = Inches(1.2)
    _rect(slide, (W - div_w) / 2, Inches(4.3), div_w, Inches(0.05), colors["accent"])

    # Summary
    _textbox(slide, s.get("summary", ""), MARGIN, Inches(4.5), W - MARGIN * 2, Inches(1.8),
             18, bold=False, italic=True, color=colors["muted"], align=PP_ALIGN.CENTER)

    # Bottom bar
    _rect(slide, 0, H - Inches(0.06), W, Inches(0.06), colors["accent2"])

    # Slide number
    num_str = str(idx).zfill(2)
    _textbox(slide, num_str, Inches(12), Inches(6.9), Inches(1.2), Inches(0.5),
             13, bold=True, color=colors["muted"], align=PP_ALIGN.RIGHT)


def build_pptx(slides: list[dict], theme: str) -> bytes:
    colors = THEMES.get(theme, THEMES["purple"])
    prs = Presentation()
    prs.slide_width = W
    prs.slide_height = H
    blank = prs.slide_layouts[6]

    for i, s in enumerate(slides, start=1):
        slide_type = s.get("type", "")
        if slide_type not in ("cover", "content", "conclusion"):
            continue
        slide = prs.slides.add_slide(blank)
        if slide_type == "cover":
            _cover_slide(slide, s, colors)
        elif slide_type == "content":
            _content_slide(slide, s, colors, i)
        elif slide_type == "conclusion":
            _conclusion_slide(slide, s, colors, i)

    buf = io.BytesIO()
    prs.save(buf)
    return buf.getvalue()
